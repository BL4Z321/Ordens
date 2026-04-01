from django.shortcuts import redirect, render
from django.contrib import messages
from django.contrib.auth.decorators import login_required
from django.contrib.auth import logout, authenticate, login
from django.contrib.messages import constants
from django.http import HttpResponse

from estoque.models import ModeloInsumo
from expedicao.models import Expedicao, ExpedicaoEnum
from insumos.models import Insumo
from modelos_customizados.models import Modelo
from ordens_producao.models import OrdemProducao, StatusOPEnum
from produtos.models import Produto
from usuarios.models import RoleEnum, Usuario

import io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
import json

def login_view(request):
    if request.method == 'GET':
        return render(request, 'login.html')
    
    elif request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        user = authenticate(request, username=username, password=password)
        if user is not None:
            try:
                usuario = Usuario.objects.get(user=user)
            except Usuario.DoesNotExist:
                messages.add_message(request, constants.ERROR, 'Usuario não cadastrado.')
                return redirect('login')
            
            if usuario.ativo:
                login(request, user)
                messages.add_message(request, constants.SUCCESS, f'Bem-vindo, {user.username}')
                
                if usuario.role == RoleEnum.ADMIN:
                    if not user.is_staff:
                        user.is_staff = True
                        user.is_superuser = True
                        user.save()

                    return redirect('/admin/')
                
                elif usuario.role == RoleEnum.GESTOR:
                    return redirect('dashboard')
                
                elif usuario.role == RoleEnum.OPERADOR:
                    return redirect('dash_operador')
                
                elif usuario.role == RoleEnum.VIEWER:
                    return redirect('viewer')
                
            else:
                messages.add_message(request, constants.ERROR, 'Usuário inativo. Contate o administrador.')
                return redirect('login')
        else:
            messages.add_message(request, constants.ERROR, 'Usuário ou senha inváidos.')
            return redirect('login')
        
    return render(request, 'login.html')
    
@login_required
def sair(request):
    logout(request)
    messages.add_message(request, constants.INFO, 'Você saiu do sistema!')
    return redirect('login')

@login_required
def dashboard_gestor(request):
    usuario = Usuario.objects.get(user=request.user)

    if usuario.role != RoleEnum.GESTOR:
        return render(request, 'login.html', {'mensagem': 'Acesso restrito ao gestor.'})

    total_ordens = OrdemProducao.objects.count()
    ordens_ativas = OrdemProducao.objects.filter(status=StatusOPEnum.EM_PRODUCAO).count()
    ordens_concluidas = OrdemProducao.objects.filter(status=StatusOPEnum.CONCLUIDA).count()
    ordens_pendentes = OrdemProducao.objects.filter(status=StatusOPEnum.PENDENTE).count()
    ordens_canceladas = OrdemProducao.objects.filter(status=StatusOPEnum.CANCELADA).count()
    ordens_bloqueadas = OrdemProducao.objects.filter(status=StatusOPEnum.BLOQUEADA).count()

    produtos_disponiveis = Produto.objects.count()

    modelos_custom = Modelo.objects.count()

    insumos_baixos = Insumo.objects.filter(estoque_atual__lt=5).count()

    expedicao = Expedicao.objects.count()
    expedi_transporte = Expedicao.objects.filter(status=ExpedicaoEnum.EM_TRANSPORTE).count()
    expedi_entregue = Expedicao.objects.filter(status=ExpedicaoEnum.ENTREGUE).count()
    expedi_retornado = Expedicao.objects.filter(status=ExpedicaoEnum.RETORNADO).count()
    expedi_enviado = Expedicao.objects.filter(status=ExpedicaoEnum.ENVIADO).count()
    expedi_pendente = Expedicao.objects.filter(status=ExpedicaoEnum.PENDENTE).count()

    estoque = ModeloInsumo.objects.count()

    contexto = {
        'usuario': usuario,
        'total_ordens': total_ordens,
        'ordens_ativas': ordens_ativas,
        'ordens_concluidas': ordens_concluidas,
        'ordens_pendentes': ordens_pendentes,
        'ordens_canceladas': ordens_canceladas,
        'ordens_bloqueadas': ordens_bloqueadas,
        'produtos_disponiveis': produtos_disponiveis,
        'modelos_custom': modelos_custom,
        'insumos_baixos': insumos_baixos,
        'expedicao': expedicao,
        'expedi_transporte': expedi_transporte,
        'expedi_entregue': expedi_entregue,
        'expedi_retornado': expedi_retornado,
        'expedi_enviado': expedi_enviado,
        'expedi_pendente': expedi_pendente,
        'estoque': estoque,
    }

    return render(request, 'gestor/dashboard.html', contexto)

@login_required
def dashboard_operador(request):
    usuario = Usuario.objects.get(user=request.user)

    if usuario.role != RoleEnum.OPERADOR:
        return render(request, 'login.html', {'mensagem': 'Acesso restrito ao gestor.'})

    total_ordens = OrdemProducao.objects.count()
    ordens_ativas = OrdemProducao.objects.filter(status=StatusOPEnum.EM_PRODUCAO).count()
    ordens_concluidas = OrdemProducao.objects.filter(status=StatusOPEnum.CONCLUIDA).count()
    ordens_pendentes = OrdemProducao.objects.filter(status=StatusOPEnum.PENDENTE).count()
    ordens_canceladas = OrdemProducao.objects.filter(status=StatusOPEnum.CANCELADA).count()
    ordens_bloqueadas = OrdemProducao.objects.filter(status=StatusOPEnum.BLOQUEADA).count()

    contexto = {
        'usuario': usuario,
        'total_ordens': total_ordens,
        'ordens_ativas': ordens_ativas,
        'ordens_concluidas': ordens_concluidas,
        'ordens_pendentes': ordens_pendentes,
        'ordens_canceladas': ordens_canceladas,
        'ordens_bloqueadas': ordens_bloqueadas,
    }

    return render(request, 'operador/dashboard.html', contexto)

@login_required
def dashboard_viewer(request):
    usuario = Usuario.objects.get(user=request.user)

    if usuario.role != RoleEnum.VIEWER:
        return render(request, 'login.html', {'mensagem': 'Acesso restrito ao gestor.'})

    total_ordens = OrdemProducao.objects.count()
    ordens_ativas = OrdemProducao.objects.filter(status=StatusOPEnum.EM_PRODUCAO).count()
    ordens_concluidas = OrdemProducao.objects.filter(status=StatusOPEnum.CONCLUIDA).count()
    ordens_pendentes = OrdemProducao.objects.filter(status=StatusOPEnum.PENDENTE).count()
    ordens_canceladas = OrdemProducao.objects.filter(status=StatusOPEnum.CANCELADA).count()
    ordens_bloqueadas = OrdemProducao.objects.filter(status=StatusOPEnum.BLOQUEADA).count()

    # métricas adicionais do sistema para dashboards/relatório
    produtos_disponiveis = Produto.objects.count()
    modelos_custom = Modelo.objects.count()
    insumos_baixos = Insumo.objects.filter(estoque_atual__lt=5).count()
    expedicao = Expedicao.objects.count()
    expedi_transporte = Expedicao.objects.filter(status=ExpedicaoEnum.EM_TRANSPORTE).count()
    expedi_entregue = Expedicao.objects.filter(status=ExpedicaoEnum.ENTREGUE).count()
    expedi_retornado = Expedicao.objects.filter(status=ExpedicaoEnum.RETORNADO).count()
    expedi_enviado = Expedicao.objects.filter(status=ExpedicaoEnum.ENVIADO).count()
    expedi_pendente = Expedicao.objects.filter(status=ExpedicaoEnum.PENDENTE).count()
    estoque = ModeloInsumo.objects.count()

    metrics = {
        'total_ordens': total_ordens,
        'ordens_ativas': ordens_ativas,
        'ordens_concluidas': ordens_concluidas,
        'ordens_pendentes': ordens_pendentes,
        'ordens_canceladas': ordens_canceladas,
        'ordens_bloqueadas': ordens_bloqueadas,
        'produtos_disponiveis': produtos_disponiveis,
        'modelos_custom': modelos_custom,
        'insumos_baixos': insumos_baixos,
        'expedicao': expedicao,
        'expedi_transporte': expedi_transporte,
        'expedi_entregue': expedi_entregue,
        'expedi_retornado': expedi_retornado,
        'expedi_enviado': expedi_enviado,
        'expedi_pendente': expedi_pendente,
        'estoque': estoque,
    }

    contexto = {
        'usuario': usuario,
        'metrics_json': json.dumps(metrics),
        'metrics': metrics,
    }

    return render(request, 'viewer/dashboard.html', contexto)

def _create_bar_chart_bytes(metrics: dict, title: str = "Métricas") -> bytes:
    labels = list(metrics.keys())
    values = list(metrics.values())
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.bar(labels, values, color='tab:blue')
    ax.set_title(title)
    ax.set_xticklabels(labels, rotation=45, ha='right', fontsize=8)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='png')
    plt.close(fig)
    buf.seek(0)
    return buf.read()

def _create_pie_chart_bytes(metrics: dict, title: str = "Distribuição") -> bytes:
    labels = list(metrics.keys())
    values = list(metrics.values())
    fig, ax = plt.subplots(figsize=(6, 6))
    ax.pie(values, labels=labels, autopct='%1.1f%%', startangle=90)
    ax.set_title(title)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='png')
    plt.close(fig)
    buf.seek(0)
    return buf.read()

@login_required
def viewer_export_pptx(request):
    usuario = Usuario.objects.get(user=request.user)
    if usuario.role != RoleEnum.VIEWER:
        return render(request, 'login.html', {'mensagem': 'Acesso restrito ao viewer.'})
    # coletar métricas (mesmas da dashboard) com chaves consistentes
    metrics_map = {
        'total_ordens': OrdemProducao.objects.count(),
        'ordens_ativas': OrdemProducao.objects.filter(status=StatusOPEnum.EM_PRODUCAO).count(),
        'ordens_concluidas': OrdemProducao.objects.filter(status=StatusOPEnum.CONCLUIDA).count(),
        'ordens_pendentes': OrdemProducao.objects.filter(status=StatusOPEnum.PENDENTE).count(),
        'ordens_canceladas': OrdemProducao.objects.filter(status=StatusOPEnum.CANCELADA).count(),
        'ordens_bloqueadas': OrdemProducao.objects.filter(status=StatusOPEnum.BLOQUEADA).count(),
        'produtos_disponiveis': Produto.objects.count(),
        'modelos_custom': Modelo.objects.count(),
        'insumos_baixos': Insumo.objects.filter(estoque_atual__lt=5).count(),
        'expedicao': Expedicao.objects.count(),
        'estoque_modelos_insumo': ModeloInsumo.objects.count(),
    }

    # rótulos legíveis
    labels_map = {
        'total_ordens': 'Total Ordens',
        'ordens_ativas': 'Ordens Em Produção',
        'ordens_concluidas': 'Ordens Concluídas',
        'ordens_pendentes': 'Ordens Pendentes',
        'ordens_canceladas': 'Ordens Canceladas',
        'ordens_bloqueadas': 'Ordens Bloqueadas',
        'produtos_disponiveis': 'Produtos Disponíveis',
        'modelos_custom': 'Modelos Customizados',
        'insumos_baixos': 'Insumos Baixos (<5)',
        'expedicao': 'Expedições',
        'estoque_modelos_insumo': 'Estoque (Modelos de Insumo)',
    }

    # parâmetros de filtro do frontend
    selected = request.GET.get('selected')
    chart_type = request.GET.get('chart', 'both')  # 'bar', 'pie', 'both'
    if selected:
        keys = [k for k in selected.split(',') if k in metrics_map]
    else:
        # padrão: exibir todas
        keys = list(metrics_map.keys())

    chosen = {labels_map[k]: metrics_map[k] for k in keys}

    prs = Presentation()

    # título
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Relatório do Sistema"
    subtitle = slide.placeholders[1]
    subtitle.text = f"Gerado por: {usuario.user.username}"

    # slide de métricas (texto)
    body_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(body_slide_layout)
    slide.shapes.title.text = "Métricas Selecionadas"
    tf = slide.shapes.placeholders[1].text_frame
    for k, v in chosen.items():
        p = tf.add_paragraph()
        p.text = f"{k}: {v}"
        p.level = 1

    # gráfico de barras (se solicitado)
    if chart_type in ('bar', 'both'):
        chart_bytes = _create_bar_chart_bytes(chosen, title="Visão Geral - Barras")
        img_path = io.BytesIO(chart_bytes)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Gráfico de Barras"
        left = Inches(0.6)
        top = Inches(1.5)
        slide.shapes.add_picture(img_path, left, top, width=Inches(8))

    # gráfico de pizza (se solicitado)
    if chart_type in ('pie', 'both'):
        pie_bytes = _create_pie_chart_bytes(chosen, title="Distribuição")
        img_path = io.BytesIO(pie_bytes)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Gráfico de Pizza"
        left = Inches(2)
        top = Inches(1.5)
        slide.shapes.add_picture(img_path, left, top, width=Inches(6))

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)

    response = HttpResponse(out.read(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = 'attachment; filename="relatorio_sistema.pptx"'
    return response